import json
from pathlib import Path
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from backend_core.storage import export_pptx_path

# ── Color Palette ─────────────────────────────────────
DARK       = RGBColor(0x0f, 0x17, 0x2a)
DARK2      = RGBColor(0x1e, 0x29, 0x3b)
PRIMARY    = RGBColor(0x4f, 0x46, 0xe5)
CYAN       = RGBColor(0x06, 0xb6, 0xd4)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GRAY       = RGBColor(0x94, 0xa3, 0xb8)
LIGHT      = RGBColor(0xf1, 0xf5, 0xf9)
SUCCESS    = RGBColor(0x05, 0x96, 0x69)
WARNING    = RGBColor(0xd9, 0x77, 0x06)
DANGER     = RGBColor(0xdc, 0x26, 0x26)
PINK       = RGBColor(0xf4, 0x72, 0xb6)

W = Inches(13.33)
H = Inches(7.5)

def _load(path):
    if path and Path(path).exists():
        try:
            return json.loads(Path(path).read_text(encoding="utf-8"))
        except:
            return {}
    return {}

def _slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def _bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def _rect(slide, l, t, w, h, color, alpha=None):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def _txt(slide, text, l, t, w, h, size=18, bold=False, color=WHITE,
         align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb

def _metric_box(slide, l, t, w, h, label, value, label_color=GRAY,
                value_color=CYAN, bg=DARK2):
    _rect(slide, l, t, w, h, bg)
    # Accent top border
    _rect(slide, l, t, w, Inches(0.05), PRIMARY)
    # Label
    _txt(slide, label, l+Inches(0.15), t+Inches(0.15), w-Inches(0.3),
         Inches(0.3), size=10, bold=False, color=label_color, align=PP_ALIGN.LEFT)
    # Value
    val_size = 28 if len(str(value)) <= 6 else 22 if len(str(value)) <= 10 else 16
    _txt(slide, value, l+Inches(0.15), t+Inches(0.45), w-Inches(0.3),
         Inches(0.6), size=val_size, bold=True, color=value_color, align=PP_ALIGN.LEFT)

def _accent_line(slide, l, t, w, color=PRIMARY):
    _rect(slide, l, t, w, Inches(0.04), color)

# ── Slide 1: Cover ────────────────────────────────────
def _slide_cover(prs, dataset_id):
    slide = _slide(prs)
    _bg(slide, DARK)

    # Left accent bar
    _rect(slide, Inches(0), Inches(0), Inches(0.08), H, PRIMARY)

    # Gradient-like background panels
    _rect(slide, Inches(0.08), Inches(0), Inches(6), H, DARK2)

    # Right decorative panel
    _rect(slide, Inches(6.08), Inches(0), W-Inches(6.08), H, DARK)

    # Decorative circles
    _rect(slide, Inches(8), Inches(-1), Inches(4), Inches(4), RGBColor(0x1e,0x1b,0x4b))
    _rect(slide, Inches(10), Inches(4), Inches(5), Inches(5), RGBColor(0x0c,0x1a,0x33))

    # Logo text
    _txt(slide, "InsightAdvisor", Inches(0.5), Inches(1.2), Inches(5.5), Inches(0.8),
         size=14, bold=True, color=CYAN, align=PP_ALIGN.LEFT)

    # Main title
    _txt(slide, "AI-Powered", Inches(0.5), Inches(1.9), Inches(5.5), Inches(0.9),
         size=44, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    _txt(slide, "Business Analytics", Inches(0.5), Inches(2.6), Inches(5.5), Inches(0.9),
         size=44, bold=True, color=PRIMARY, align=PP_ALIGN.LEFT)
    _txt(slide, "Report", Inches(0.5), Inches(3.3), Inches(5.5), Inches(0.9),
         size=44, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    # Subtitle
    _txt(slide, "End-to-End Analytics · Predictive Modeling · AI Recommendations",
         Inches(0.5), Inches(4.3), Inches(5.5), Inches(0.5),
         size=11, bold=False, color=GRAY, align=PP_ALIGN.LEFT)

    # Meta info
    now = datetime.now().strftime("%B %d, %Y")
    _txt(slide, f"Generated: {now}", Inches(0.5), Inches(5.2), Inches(5.5), Inches(0.35),
         size=10, bold=False, color=GRAY, align=PP_ALIGN.LEFT)
    _txt(slide, f"Dataset: {dataset_id[:20]}...", Inches(0.5), Inches(5.5), Inches(5.5), Inches(0.35),
         size=10, bold=False, color=GRAY, align=PP_ALIGN.LEFT)
    _txt(slide, "insightadvisor.co", Inches(0.5), Inches(5.8), Inches(5.5), Inches(0.35),
         size=10, bold=False, color=CYAN, align=PP_ALIGN.LEFT)

    # Right side: section list
    sections = [
        ("01", "Data Quality"),
        ("02", "Exploratory Analysis"),
        ("03", "Customer Segments"),
        ("04", "Sales Forecast"),
        ("05", "Regression"),
        ("06", "Association Rules"),
        ("07", "Anomaly Detection"),
        ("08", "AI Recommendations"),
    ]
    for i, (num, name) in enumerate(sections):
        y = Inches(0.6) + i * Inches(0.8)
        _rect(slide, Inches(7), y, Inches(0.5), Inches(0.5), PRIMARY)
        _txt(slide, num, Inches(7), y, Inches(0.5), Inches(0.5),
             size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _txt(slide, name, Inches(7.6), y+Inches(0.05), Inches(5), Inches(0.45),
             size=13, bold=False, color=WHITE, align=PP_ALIGN.LEFT)

# ── Slide 2: Data Quality ─────────────────────────────
def _slide_quality(prs, report):
    slide = _slide(prs)
    _bg(slide, DARK)
    _rect(slide, Inches(0), Inches(0), Inches(0.08), H, SUCCESS)

    # Header
    _rect(slide, Inches(0.08), Inches(0), W, Inches(1.2), DARK2)
    _txt(slide, "01  —  Data Quality & Cleaning", Inches(0.4), Inches(0.15),
         Inches(8), Inches(0.5), size=13, bold=False, color=GRAY)
    _txt(slide, "Data Quality Report", Inches(0.4), Inches(0.55),
         Inches(10), Inches(0.55), size=28, bold=True, color=WHITE)

    score = report.get("quality_score", 0)
    score_color = SUCCESS if score >= 80 else WARNING if score >= 60 else DANGER

    outliers = report.get("outliers_removed", 0)
    if isinstance(outliers, dict):
        outliers = sum(outliers.values())

    metrics = [
        ("QUALITY SCORE",      f"{score}/100",                   score_color),
        ("ROWS BEFORE",        f"{report.get('rows_before',0):,}", CYAN),
        ("ROWS AFTER",         f"{report.get('rows_after',0):,}",  CYAN),
        ("DUPLICATES REMOVED", str(report.get("duplicates_removed",0)), WARNING),
        ("OUTLIERS REMOVED",   str(outliers),                     WARNING),
    ]
    mw = (W - Inches(0.8)) / len(metrics)
    for i, (label, value, color) in enumerate(metrics):
        _metric_box(slide, Inches(0.4)+i*mw, Inches(1.4), mw-Inches(0.1),
                    Inches(1.5), label, value, value_color=color)

    # Missing values table
    missing = report.get("missing_before", {})
    if missing:
        _txt(slide, "Missing Values by Column", Inches(0.4), Inches(3.1),
             Inches(12), Inches(0.4), size=14, bold=True, color=WHITE)
        items = list(missing.items())[:8]
        col_w = (W - Inches(0.8)) / len(items)
        for i, (col, count) in enumerate(items):
            x = Inches(0.4) + i * col_w
            _rect(slide, x, Inches(3.6), col_w-Inches(0.05), Inches(0.8), DARK2)
            _txt(slide, col, x, Inches(3.65), col_w-Inches(0.1),
                 Inches(0.35), size=9, color=GRAY, align=PP_ALIGN.CENTER)
            c = DANGER if count > 1000 else WARNING if count > 0 else SUCCESS
            _txt(slide, str(count), x, Inches(4.0), col_w-Inches(0.1),
                 Inches(0.4), size=16, bold=True, color=c, align=PP_ALIGN.CENTER)

# ── Slide 3: EDA ──────────────────────────────────────
def _slide_eda(prs, report):
    slide = _slide(prs)
    _bg(slide, DARK)
    _rect(slide, Inches(0), Inches(0), Inches(0.08), H, CYAN)
    _rect(slide, Inches(0.08), Inches(0), W, Inches(1.2), DARK2)
    _txt(slide, "02  —  Exploratory Data Analysis", Inches(0.4), Inches(0.15),
         Inches(8), Inches(0.5), size=13, bold=False, color=GRAY)
    _txt(slide, "Dataset Overview", Inches(0.4), Inches(0.55),
         Inches(10), Inches(0.55), size=28, bold=True, color=WHITE)

    stats = report.get("summary_statistics", {}) or report.get("numeric_summary", {})
    top_cats = report.get("top_categories", {})

    # Stats table
    if stats:
        _txt(slide, "Summary Statistics", Inches(0.4), Inches(1.4),
             Inches(12), Inches(0.4), size=14, bold=True, color=WHITE)
        headers = ["Column", "Mean", "Std", "Min", "Max"]
        col_widths = [Inches(2.5), Inches(1.8), Inches(1.8), Inches(1.8), Inches(1.8)]
        # Header row
        x = Inches(0.4)
        for i, h in enumerate(headers):
            _rect(slide, x, Inches(1.9), col_widths[i]-Inches(0.02), Inches(0.4), PRIMARY)
            _txt(slide, h, x, Inches(1.95), col_widths[i],
                 Inches(0.35), size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
            x += col_widths[i]
        # Data rows
        for ri, (col, s) in enumerate(list(stats.items())[:5]):
            y = Inches(2.3) + ri * Inches(0.45)
            bg = DARK2 if ri % 2 == 0 else DARK
            x = Inches(0.4)
            vals = [col, round(s.get("mean",0),2), round(s.get("std",0),2),
                    round(s.get("min",0),2), round(s.get("max",0),2)]
            for ci, val in enumerate(vals):
                _rect(slide, x, y, col_widths[ci]-Inches(0.02), Inches(0.42), bg)
                _txt(slide, str(val), x+Inches(0.05), y+Inches(0.05),
                     col_widths[ci]-Inches(0.1), Inches(0.35),
                     size=10, color=CYAN if ci==0 else WHITE, align=PP_ALIGN.CENTER)
                x += col_widths[ci]

    # Top products
    if top_cats:
        col_key = next(iter(top_cats))
        items = top_cats[col_key]
        if isinstance(items, dict):
            vals = list(zip(items.get("values",[])[:6], items.get("counts",[])[:6]))
        else:
            vals = [(i.get("value",""), i.get("count",0)) for i in items[:6]]

        _txt(slide, f"Top Values — {col_key}", Inches(0.4), Inches(4.7),
             Inches(12), Inches(0.4), size=14, bold=True, color=WHITE)
        max_count = max((v[1] for v in vals), default=1)
        bar_w = (W - Inches(0.8)) / len(vals) if vals else Inches(1)
        for i, (name, count) in enumerate(vals):
            x = Inches(0.4) + i * bar_w
            bar_h = Inches(1.2) * (count / max_count) if max_count else Inches(0.1)
            _rect(slide, x+Inches(0.05), Inches(6.8)-bar_h, bar_w-Inches(0.1), bar_h, PRIMARY)
            _txt(slide, str(count), x, Inches(6.85), bar_w,
                 Inches(0.3), size=9, bold=True, color=CYAN, align=PP_ALIGN.CENTER)
            short = str(name)[:12]
            _txt(slide, short, x, Inches(7.15), bar_w,
                 Inches(0.3), size=8, color=GRAY, align=PP_ALIGN.CENTER)

# ── Slide 4: Clustering ───────────────────────────────
def _slide_clustering(prs, report):
    slide = _slide(prs)
    _bg(slide, DARK)
    _rect(slide, Inches(0), Inches(0), Inches(0.08), H, PINK)
    _rect(slide, Inches(0.08), Inches(0), W, Inches(1.2), DARK2)
    _txt(slide, "03  —  Customer Segmentation", Inches(0.4), Inches(0.15),
         Inches(8), Inches(0.5), size=13, bold=False, color=GRAY)
    _txt(slide, "KMeans + RFM Analysis", Inches(0.4), Inches(0.55),
         Inches(10), Inches(0.55), size=28, bold=True, color=WHITE)

    total = sum(p.get("size",0) for p in report.get("cluster_profiles",[]))
    metrics = [
        ("CLUSTERS", str(report.get("k","N/A")), CYAN),
        ("SILHOUETTE", str(round(report.get("silhouette_score",0),4)), SUCCESS),
        ("DAVIES-BOULDIN", str(round(report.get("davies_bouldin",0),4)), WARNING),
        ("TOTAL CUSTOMERS", f"{total:,}", PINK),
    ]
    mw = (W - Inches(0.8)) / 4
    for i, (label, value, color) in enumerate(metrics):
        _metric_box(slide, Inches(0.4)+i*mw, Inches(1.4), mw-Inches(0.1),
                    Inches(1.3), label, value, value_color=color)

    profiles = report.get("cluster_profiles", [])
    seg_colors = [PRIMARY, CYAN, PINK, SUCCESS, WARNING, DANGER]
    if profiles:
        _txt(slide, "Customer Segments", Inches(0.4), Inches(2.9),
             Inches(12), Inches(0.4), size=14, bold=True, color=WHITE)
        cw = (W - Inches(0.8)) / min(len(profiles), 4)
        for i, p in enumerate(profiles[:4]):
            x = Inches(0.4) + i * cw
            color = seg_colors[i % len(seg_colors)]
            _rect(slide, x, Inches(3.4), cw-Inches(0.1), Inches(3.8), DARK2)
            _rect(slide, x, Inches(3.4), cw-Inches(0.1), Inches(0.06), color)
            name = p.get("name", f"Cluster {p.get('cluster',i)}")
            _txt(slide, name, x+Inches(0.1), Inches(3.5), cw-Inches(0.2),
                 Inches(0.5), size=11, bold=True, color=color)
            _txt(slide, f"{p.get('size',0):,} customers", x+Inches(0.1), Inches(3.95),
                 cw-Inches(0.2), Inches(0.35), size=10, color=GRAY)
            rfm = [
                ("Recency",   round(p.get("Recency",   p.get("avg_recency",  0)),1), "d"),
                ("Frequency", round(p.get("Frequency", p.get("avg_frequency",0)),1), "x"),
                ("Monetary",  round(p.get("Monetary",  p.get("avg_monetary", 0)),2), ""),
            ]
            for j, (label, val, unit) in enumerate(rfm):
                y = Inches(4.4) + j * Inches(0.85)
                _txt(slide, label, x+Inches(0.1), y, cw-Inches(0.2),
                     Inches(0.3), size=9, color=GRAY)
                _txt(slide, f"{val}{unit}", x+Inches(0.1), y+Inches(0.28),
                     cw-Inches(0.2), Inches(0.4), size=14, bold=True, color=color)

# ── Slide 5: Forecasting ──────────────────────────────
def _slide_forecasting(prs, report):
    slide = _slide(prs)
    _bg(slide, DARK)
    _rect(slide, Inches(0), Inches(0), Inches(0.08), H, SUCCESS)
    _rect(slide, Inches(0.08), Inches(0), W, Inches(1.2), DARK2)
    _txt(slide, "04  —  Sales Forecasting", Inches(0.4), Inches(0.15),
         Inches(8), Inches(0.5), size=13, bold=False, color=GRAY)
    _txt(slide, "ARIMA Time Series Model", Inches(0.4), Inches(0.55),
         Inches(10), Inches(0.55), size=28, bold=True, color=WHITE)

    m = report.get("metrics", {})
    s = report.get("summary", {})
    trend = s.get("trend","stable")
    trend_color = SUCCESS if trend=="increasing" else DANGER if trend=="decreasing" else WARNING

    metrics = [
        ("MODEL",          report.get("model","ARIMA"),  CYAN),
        ("RMSE",           round(m.get("rmse",0),2),     WARNING),
        ("MAPE",           f"{round(m.get('mape',0),2)}%", WARNING),
        ("TOTAL FORECAST", f"{round(s.get('total_forecast_revenue',0),0):,.0f}", SUCCESS),
        ("TREND",          trend.upper(),                trend_color),
    ]
    mw = (W - Inches(0.8)) / 5
    for i, (label, value, color) in enumerate(metrics):
        _metric_box(slide, Inches(0.4)+i*mw, Inches(1.4), mw-Inches(0.1),
                    Inches(1.3), label, str(value), value_color=color)

    # Forecast table
    forecast = report.get("forecast", [])
    if forecast:
        _txt(slide, "30-Day Forecast Sample", Inches(0.4), Inches(2.9),
             Inches(12), Inches(0.4), size=14, bold=True, color=WHITE)
        headers = ["Date", "Forecast", "Lower (95%)", "Upper (95%)"]
        cws = [Inches(2.5), Inches(2.5), Inches(2.5), Inches(2.5)]
        x = Inches(0.4)
        for i, h in enumerate(headers):
            _rect(slide, x, Inches(3.4), cws[i]-Inches(0.02), Inches(0.4), PRIMARY)
            _txt(slide, h, x, Inches(3.45), cws[i], Inches(0.35),
                 size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
            x += cws[i]
        for ri, f in enumerate(forecast[:7]):
            y = Inches(3.85) + ri * Inches(0.45)
            bg = DARK2 if ri%2==0 else DARK
            x = Inches(0.4)
            vals = [str(f.get("date",""))[:10], round(f.get("forecast",0),2),
                    round(f.get("lower_bound",0),2), round(f.get("upper_bound",0),2)]
            for ci, val in enumerate(vals):
                _rect(slide, x, y, cws[ci]-Inches(0.02), Inches(0.42), bg)
                _txt(slide, str(val), x, y+Inches(0.05), cws[ci],
                     Inches(0.35), size=10,
                     color=CYAN if ci==1 else WHITE, align=PP_ALIGN.CENTER)
                x += cws[ci]

# ── Slide 6: Association Rules ────────────────────────
def _slide_apriori(prs, report):
    slide = _slide(prs)
    _bg(slide, DARK)
    _rect(slide, Inches(0), Inches(0), Inches(0.08), H, WARNING)
    _rect(slide, Inches(0.08), Inches(0), W, Inches(1.2), DARK2)
    _txt(slide, "06  —  Association Rules", Inches(0.4), Inches(0.15),
         Inches(8), Inches(0.5), size=13, bold=False, color=GRAY)
    _txt(slide, "Market Basket Analysis (Apriori)", Inches(0.4), Inches(0.55),
         Inches(10), Inches(0.55), size=28, bold=True, color=WHITE)

    s = report.get("summary", {})
    rules = report.get("top_rules", [])
    max_conf = s.get("max_confidence") or (rules[0].get("confidence",0) if rules else 0)
    max_lift = s.get("max_lift") or (rules[0].get("lift",0) if rules else 0)

    metrics = [
        ("TOTAL RULES",    str(report.get("total_rules_found", len(rules))), CYAN),
        ("MAX CONFIDENCE", f"{round(max_conf*100,1)}%", SUCCESS),
        ("MAX LIFT",       str(round(max_lift,2)), WARNING),
    ]
    mw = (W - Inches(0.8)) / 3
    for i, (label, value, color) in enumerate(metrics):
        _metric_box(slide, Inches(0.4)+i*mw, Inches(1.4), mw-Inches(0.1),
                    Inches(1.3), label, value, value_color=color)

    if rules:
        _txt(slide, "Top Rules in Plain English", Inches(0.4), Inches(2.9),
             Inches(12), Inches(0.4), size=14, bold=True, color=WHITE)
        for i, rule in enumerate(rules[:5]):
            y = Inches(3.4) + i * Inches(0.75)
            _rect(slide, Inches(0.4), y, W-Inches(0.8), Inches(0.68), DARK2)
            _rect(slide, Inches(0.4), y, Inches(0.06), Inches(0.68), WARNING)
            eng = rule.get("english", f"If {rule.get('antecedents','')} then {rule.get('consequents','')}")
            _txt(slide, f"{i+1}.  {eng}", Inches(0.6), y+Inches(0.1),
                 Inches(10), Inches(0.5), size=10, color=WHITE)
            conf = f"{round(rule.get('confidence',0)*100,1)}%"
            lift = f"{round(rule.get('lift',0),2)}x"
            _txt(slide, f"Conf: {conf}  |  Lift: {lift}", Inches(11), y+Inches(0.2),
                 Inches(2.1), Inches(0.35), size=9, bold=True, color=WARNING, align=PP_ALIGN.RIGHT)

# ── Slide 7: Anomaly ──────────────────────────────────
def _slide_anomaly(prs, report):
    slide = _slide(prs)
    _bg(slide, DARK)
    _rect(slide, Inches(0), Inches(0), Inches(0.08), H, DANGER)
    _rect(slide, Inches(0.08), Inches(0), W, Inches(1.2), DARK2)
    _txt(slide, "07  —  Anomaly Detection", Inches(0.4), Inches(0.15),
         Inches(8), Inches(0.5), size=13, bold=False, color=GRAY)
    _txt(slide, "Isolation Forest Algorithm", Inches(0.4), Inches(0.55),
         Inches(10), Inches(0.55), size=28, bold=True, color=WHITE)

    summ = report.get("summary", {})
    metrics = [
        ("TOTAL ROWS",      f"{report.get('total_rows',0):,}", CYAN),
        ("ANOMALIES",       f"{report.get('total_anomalies',0):,}", DANGER),
        ("ANOMALY %",       f"{report.get('anomaly_pct',0)}%", WARNING),
        ("HIGH SEVERITY",   str(summ.get("high_severity",0)), DANGER),
        ("MED SEVERITY",    str(summ.get("medium_severity",0)), WARNING),
    ]
    mw = (W - Inches(0.8)) / 5
    for i, (label, value, color) in enumerate(metrics):
        _metric_box(slide, Inches(0.4)+i*mw, Inches(1.4), mw-Inches(0.1),
                    Inches(1.3), label, value, value_color=color)

    top = report.get("top_anomalies", [])
    if top:
        _txt(slide, "Top Anomalies Detected", Inches(0.4), Inches(2.9),
             Inches(12), Inches(0.4), size=14, bold=True, color=WHITE)
        for i, a in enumerate(top[:5]):
            y = Inches(3.4) + i * Inches(0.75)
            sev = a.get("severity","Low")
            sev_color = DANGER if sev=="High" else WARNING if sev=="Medium" else GRAY
            _rect(slide, Inches(0.4), y, W-Inches(0.8), Inches(0.68), DARK2)
            _rect(slide, Inches(0.4), y, Inches(0.06), Inches(0.68), sev_color)
            _txt(slide, f"Row {a.get('row_index','?')}  [{sev}]",
                 Inches(0.6), y+Inches(0.05), Inches(3), Inches(0.3),
                 size=10, bold=True, color=sev_color)
            _txt(slide, a.get("explanation",""), Inches(0.6), y+Inches(0.32),
                 Inches(12), Inches(0.35), size=9, color=GRAY)

# ── Slide 8: AI Recommendations ──────────────────────
def _slide_ai(prs, report):
    slide = _slide(prs)
    _bg(slide, DARK)
    _rect(slide, Inches(0), Inches(0), Inches(0.08), H, PRIMARY)
    _rect(slide, Inches(0.08), Inches(0), W, Inches(1.2), DARK2)
    _txt(slide, "08  —  AI Recommendations", Inches(0.4), Inches(0.15),
         Inches(8), Inches(0.5), size=13, bold=False, color=GRAY)
    _txt(slide, "Powered by LLM Intelligence", Inches(0.4), Inches(0.55),
         Inches(10), Inches(0.55), size=28, bold=True, color=WHITE)

    recs = report.get("recommendations", {})
    if not recs:
        _txt(slide, "Run AI analysis to generate recommendations",
             Inches(0.4), Inches(3.5), Inches(12), Inches(0.5),
             size=16, color=GRAY, align=PP_ALIGN.CENTER)
        return

    overall = recs.get("overall_health","")
    if overall:
        _rect(slide, Inches(0.4), Inches(1.35), W-Inches(0.8), Inches(0.7), DARK2)
        _rect(slide, Inches(0.4), Inches(1.35), Inches(0.06), Inches(0.7), CYAN)
        _txt(slide, overall, Inches(0.6), Inches(1.45), W-Inches(1.2), Inches(0.55),
             size=11, color=WHITE)

    priorities = [
        ("priority_1", "TODAY",      DANGER),
        ("priority_2", "THIS WEEK",  WARNING),
        ("priority_3", "THIS MONTH", SUCCESS),
    ]
    pw = (W - Inches(0.8)) / 3
    for i, (key, label, color) in enumerate(priorities):
        p = recs.get(key, {})
        if not p:
            continue
        x = Inches(0.4) + i * pw
        y = Inches(2.2)
        _rect(slide, x, y, pw-Inches(0.1), Inches(4.8), DARK2)
        _rect(slide, x, y, pw-Inches(0.1), Inches(0.5), color)
        _txt(slide, label, x+Inches(0.1), y+Inches(0.08), pw-Inches(0.2),
             Inches(0.35), size=11, bold=True, color=WHITE)
        _txt(slide, p.get("title",""), x+Inches(0.1), y+Inches(0.6),
             pw-Inches(0.2), Inches(0.5), size=11, bold=True, color=color)
        _txt(slide, p.get("action",""), x+Inches(0.1), y+Inches(1.15),
             pw-Inches(0.2), Inches(1.8), size=10, color=WHITE)
        _txt(slide, p.get("expected_impact",""), x+Inches(0.1), y+Inches(3.8),
             pw-Inches(0.2), Inches(0.6), size=10, bold=True, color=color)

# ── Main Entry Point ──────────────────────────────────
def generate_pptx(dataset_id: str) -> Path:
    out_path = export_pptx_path(dataset_id)
    base = out_path.parent.parent / "processed"

    def load(name):
        p = base / f"{dataset_id}__{name}.json"
        return _load(p) if p.exists() else {}

    cleaning   = load("cleaning_report")
    eda        = load("eda_report")
    kmeans     = load("kmeans_report")
    arima      = load("arima_report")
    apriori    = load("apriori_report")
    anomaly    = load("anomaly_report")
    ai         = load("ai_report")

    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    _slide_cover(prs, dataset_id)
    if cleaning:   _slide_quality(prs, cleaning)
    if eda:        _slide_eda(prs, eda)
    if kmeans:     _slide_clustering(prs, kmeans)
    if arima:      _slide_forecasting(prs, arima)
    if apriori:    _slide_apriori(prs, apriori)
    if anomaly:    _slide_anomaly(prs, anomaly)
    if ai:         _slide_ai(prs, ai)

    prs.save(str(out_path))
    return out_path
